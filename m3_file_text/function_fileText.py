#    function_fileText
#    Copyright (C) 2024  Diana Martin
#
#    This program is free software: you can redistribute it and/or modify
#    it under the terms of the GNU General Public License as published by
#    the Free Software Foundation, either version 3 of the License, or
#    (at your option) any later version.
#
#    This program is distributed in the hope that it will be useful,
#    but WITHOUT ANY WARRANTY; without even the implied warranty of
#    MERCHANTABILITY or FITNESS FOR A PARTICULAR PURPOSE.  See the
#    GNU General Public License for more details.

#    You should have received a copy of the GNU General Public License
#    along with this program.  If not, see <https://www.gnu.org/licenses/>.

"""
This script serves as the functional aspect of this module.
This script extracts the textual content of various file types and saves it into a plain text file. This program uses several libraries to extract text from various file types.

Usage: $ python function_fileText.py <inputFile> [-o <outputFile>]

Example: $ python function_fileText.py /path/document.docx -o /path/output.txt
"""
# Import necessary libraries.
import argparse
import os
import json
from pptx import Presentation
from docx import Document
from bs4 import BeautifulSoup
import ebooklib
from ebooklib import epub

def isEpub(filePath):
    """
    This function verifies if the input file is an EPUB file.
    
    Parameters: 
        * filePath (str): The path of the input file.

    Returns: 
        * bool: True if the file is an '.epub', False if the file is not.  
    """
    # Check if the file's extension is an '.epub'.
    return filePath.endswith('.epub')

def isHtml(filePath):
    """
    This function verifies if the input file is an HTML or HTM file.
    
    Parameters: 
        * filePath (str): The path of the input file.

    Returns: 
        * bool: True if the file is an '.html' or '.htm', False if the file is not.  
    """
    # Check if the file's extension is an '.html' or '.htm'.
    return filePath.endswith('.html') or filePath.endswith('.htm')

def isJson(filePath):
    """
    This function verifies if the input file is a JSON file.
    
    Parameters: 
        * filePath (str): The path of the input file.

    Returns: 
        * bool: True if the file is a '.json', False if the file is not.  
    """
    # Check if the file's extension is a '.json'.
    return filePath.endswith('.json')

def isPptx(filePath):
    """
    This function verifies if the input file is a PPTX file.
    
    Parameters: 
        * filePath (str): The path of the input file.

    Returns: 
        * bool: True if the file is a '.pptx', False if the file is not.  
    """
    # Check if the file's extension is a '.pptx'.
    return filePath.endswith('.pptx')

def isDocx(filePath):
    """
    This function verifies if the input file is a DOCX file.
    
    Parameters: 
        * filePath (str): The path of the input file.

    Returns: 
        * bool: True if the file is a '.docx', False if the file is not.  
    """
    # Check if the file's extension is a '.docx'.
    return filePath.endswith('.docx')

def extractTextFromEpub(filePath):
    """
    This function extracts the textual content of an EPUB file.
    
    Parameters: 
        * filePath (str): The path of the input file.

    Returns: 
        * str: The extracted content of the '.epub' file.
    """
    # Read the file using the epub.read_epub function.
    book = epub.read_epub(filePath)
    # Initialize an empty list to store the content.
    texts = []
    # Iterate through each item.
    for item in book.get_items():
        # Check if the item is a document.
        if item.get_type() == ebooklib.ITEM_DOCUMENT:
            # Decode the content of the document and append it to the text list.
            texts.append(item.get_body_content().decode('utf-8'))
    # Join the text with newline characters and retun the extracted text.
    return '\n'.join(texts)

def extractTextFromHtml(filePath):
    """
    This function extracts the textual content of an HTML or HTM file.
    
    Parameters: 
        * filePath (str): The path of the input file.

    Returns: 
        * str: The extracted content of the '.html' or '.htm' file.
    """
    # Open the input file in reading mode with 'utf-8' encoding.
    with open(filePath, 'r', encoding='utf-8') as file:
        # Parse the file content using BeautifulSoup with the 'html.parser' parser.
        soup = BeautifulSoup(file, 'html.parser')
        # Extract and return the extracted text.
        return soup.get_text()

def extractTextFromJson(filePath):
    """
    This function extracts the textual content of a JSON file.
    
    Parameters: 
        * filePath (str): The path of the input file.

    Returns: 
        * str: The extracted content of the '.json' file.
    """
    # Open the input file in reading mode with 'utf-8' encoding. 
    with open(filePath, 'r', encoding='utf-8') as file:
        # Load the content into a Python dictionary or list.
        data = json.load(file)
        # Convert the loaded content into a script and return it.
        return str(data)

def extractTextFromPptx(filePath):
    """
    This function extracts the textual content of a PPTX file.
    
    Parameters: 
        * filePath (str): The path of the input file.

    Returns: 
        * str: The extracted content of the '.pptx' file.
    """
    # Open the input file.
    prs = Presentation(filePath)
    # Initialize an empty list to store the content of the input file.
    textRuns = []
    # Iterate through each slide.
    for slide in prs.slides:
        # Iterate throught each shape in the slide.
        for shape in slide.shapes:
            # Verify if the the shape object is a text.
            if hasattr(shape, 'text'):
                # Append the text from the shape an into the text_runs list. 
                textRuns.append(shape.text)
    return '\n'.join(textRuns)

def extractTextFromDocx(filePath):
    """
    This function extracts the textual content of a DOCX file.
    
    Parameters: 
        * filePath (str): The path of the input file.

    Returns: 
        * str: The extracted content of the '.docx' file.
    """
    # Open the input file.
    doc = Document(filePath)
    # Initialize an empty list to store the content of the input file.
    fullText = []
    # Iterate through each paragraph.
    for para in doc.paragraphs:
        # Append the text from the paragraphs an into the full_text list. 
        fullText.append(para.text)
    # Join the text with newline characters and retun the extracted text.
    return '\n'.join(fullText)

def extractText(filePath):
    """
    This function determines the file format based on the file's extension.

    Parameters:
        * filePath (str): Path of the input file.
    
    Returns:
        * str: The extracted content from the files.

    Raises:
        * ValueError: In case the file format entered is not supported.
    """
    # Verifies if the file supported and calls the appropiate function.
    if isEpub(filePath):
        return extractTextFromEpub(filePath)
    elif isHtml(filePath):
        return extractTextFromHtml(filePath)
    elif isJson(filePath):
        return extractTextFromJson(filePath)
    elif isPptx(filePath):
        return extractTextFromPptx(filePath)
    elif isDocx(filePath):
        return extractTextFromDocx(filePath)
    else:
        raise ValueError(f'Unsupported file format for {filePath}')

def saveText(text, outputFile):
    """
    This function saves extracted text to a specified output file.

    Parameters:
        * text (str): Text to be saved.
        * outputFile (str): Path of the output file.

    Returns:
        * None
    """
    # Open the source file and write the content onto the output file.
    with open(outputFile, 'w', encoding='utf-8') as file:
        file.write(text)
    # Print the sucess message.
    print(f'Successfully saved extracted text to {outputFile}')

if __name__ == "__main__":
    """
    Main execution block.
    This block parses command-line arguments to extract the content of various file types to a plain text file.
    It requires the path to the input file and optionally the path to the output file.
    """
    # Configure argument parser for command-line usage.
    parser = argparse.ArgumentParser(description='Script to extract text from various file formats.')
    parser.add_argument('-i', '--input', dest='inputFile', help='Path to the input file', required=True)
    parser.add_argument('-o', '--output', dest='outputFile', help='Path to the output text file')

    # Parse command-line arguments.
    args = parser.parse_args()

    # Execute the extraction and saving process.
    try:
        extractedText = extractText(args.inputFile)
        # If the user provides the name of the output file in the arguments.
        if args.outputFile:
            # Extract text from the source file and save it with the choosen output filename.
            saveText(extractedText, args.outputFile)
        # If not, the program provides a default output name.
        else:
            # Get the source file's absolute path. 
            inputDir = os.path.dirname(os.path.abspath(args.inputFile))
            # Obtain the name of the file without the extension.
            baseName = os.path.splitext(os.path.basename(args.inputFile))[0]
            # Create the default name by appending the file's base name with "_output.txt"
            outputFile = os.path.join(inputDir, f'{baseName}_output.txt')
            # Extract text from the source file and save it with the default output filename.
            saveText(extractedText, outputFile)
    # Handle other types of errors and show message.
    except Exception as e:
        print(f'Error: {e}')
